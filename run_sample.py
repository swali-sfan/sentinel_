"""
End-to-end smoke test for the Sentinel IQ Document Formatter.

Run:
    python run_sample.py

Exercises every parser + the reformatter + the draft generator, writes
sample outputs into ./sample_output/, and prints a summary.
"""

from __future__ import annotations
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from src.parsers import parse
from src.reformatter import reformat
from src.draft import generate_draft
from src.render_pdf import render_pdf


def _write(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def _make_sample_inputs(out_dir: Path) -> dict[str, Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    files: dict[str, Path] = {}

    md = """# Legacy Strategy Note

This is an older internal note we want to bring into the Sentinel IQ format.

## 0. The Strategic Position
Old position paragraph.

## 1. Why This Matters
- Reason one
- Reason two

## 2. The Three Layers
- Foundation
- Operating System
- Acceleration

## 4. Core Frameworks
TBD

## 7. Closing
Contact: jane@example.com
"""
    files["sample.md"] = out_dir / "sample.md"
    files["sample.md"].write_text(md, encoding="utf-8")

    files["sample.txt"] = out_dir / "sample.txt"
    files["sample.txt"].write_text("Legacy Text Memo\n\nA few paragraphs of free text about the project.\n", encoding="utf-8")

    # Build a tiny .docx
    try:
        from docx import Document
        d = Document()
        d.add_heading("Legacy Word Document", level=1)
        d.add_paragraph("Some old strategy text inside a Word file.")
        d.add_heading("Section A", level=2)
        d.add_paragraph("Body of section A.")
        tbl = d.add_table(rows=2, cols=2)
        tbl.cell(0, 0).text = "Key"
        tbl.cell(0, 1).text = "Value"
        tbl.cell(1, 0).text = "Owner"
        tbl.cell(1, 1).text = "Jane"
        files["sample.docx"] = out_dir / "sample.docx"
        d.save(files["sample.docx"])
    except Exception as e:
        print(f"  ! skipped docx: {e}")

    # Build a tiny .pptx
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Legacy Slides"
        slide.placeholders[1].text = "Bullet one\nBullet two"
        files["sample.pptx"] = out_dir / "sample.pptx"
        prs.save(files["sample.pptx"])
    except Exception as e:
        print(f"  ! skipped pptx: {e}")

    return files


def main() -> int:
    sample_dir = ROOT / "sample_input"
    out_dir = ROOT / "sample_output"
    out_dir.mkdir(parents=True, exist_ok=True)

    print("== Sentinel IQ Formatter — Smoke Test ==\n")

    print("[1/4] Building sample inputs…")
    inputs = _make_sample_inputs(sample_dir)
    for name, p in inputs.items():
        print(f"   - {p.name}  ({p.stat().st_size} bytes)")

    print("\n[2/4] Reformatting each source into Sentinel IQ markdown…")
    migrated: list[Path] = []
    for idx, (name, p) in enumerate(inputs.items(), start=1):
        try:
            parsed = parse(p)
            md = reformat(
                parsed,
                sub_brand=f"SIQ-{abs(hash(name)) % 100000:05d}",
            )
            out = out_dir / f"{p.stem}.{idx:02d}.SIQ.md"
            out.write_text(md, encoding="utf-8")
            migrated.append(out)
            print(f"   ✓ {p.name:>14} → {out.name}  ({len(md):,} chars)")
        except Exception as e:
            print(f"   ✗ {p.name}: {e}")

    print("\n[3/4] Generating a DRAFT from a form…")
    draft_md = generate_draft(
        title="AI-Driven Compliance Operating System",
        vertical="Financial Services",
        owner="Jane Doe",
        sub_brand="SIQ-42001",
        contact="jane@sentinel-iq.example",
        review_cycle_days=60,
    )
    draft_path = out_dir / "AI-Driven Compliance Operating System - SIQ Draft.md"
    draft_path.write_text(draft_md, encoding="utf-8")
    print(f"   ✓ DRAFT → {draft_path.name}  ({len(draft_md):,} chars)")

    print("\n[4/4] Attempting PDF render of the draft…")
    try:
        pdf = render_pdf(draft_path)
        print(f"   ✓ PDF → {pdf}")
    except Exception as e:
        print(f"   ! PDF render skipped: {e}")
        print("     (install wkhtmltopdf or pandoc to enable this step)")

    print("\n== Done ==")
    print(f"Sample inputs : {sample_dir}")
    print(f"Sample outputs: {out_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
