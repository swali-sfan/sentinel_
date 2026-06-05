"""
Input parsers for the Sentinel IQ Document Formatter.

Each parser takes a file path and returns a ParsedDocument:
  {
    "title":  str | None,       # best-guess document title
    "body":   str,              # plain-text body, in reading order
    "tables": list[str],        # any tables, as pipe-tables
    "meta":   dict,             # format-specific metadata
    "source_format": str,       # "md" | "txt" | "docx" | "pdf" | "pptx"
  }
"""

from __future__ import annotations
from pathlib import Path
from typing import Any, Iterator


SUPPORTED_EXTS = {".md", ".txt", ".docx", ".pdf", ".pptx"}


# ---------------------------------------------------------------------------
# Markdown / Text — trivial: read, lightly normalise
# ---------------------------------------------------------------------------

def _parse_markdown_or_text(path: Path) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return {
        "title": _extract_md_title(text),
        "body":  text,
        "tables": _extract_md_tables(text),
        "meta": {"path": str(path), "size_kb": round(path.stat().st_size / 1024, 1)},
        "source_format": "md" if path.suffix.lower() == ".md" else "txt",
    }


def _extract_md_title(text: str) -> str | None:
    for line in text.splitlines():
        s = line.strip()
        if s.startswith("# "):
            return s[2:].strip()
        if s and not s.startswith("#"):
            return s
    return None


def _extract_md_tables(text: str) -> list[str]:
    tables, buf = [], []
    in_table = False
    for line in text.splitlines():
        if line.strip().startswith("|") and "|" in line.strip()[1:]:
            buf.append(line)
            in_table = True
        else:
            if in_table and buf:
                tables.append("\n".join(buf))
                buf = []
            in_table = False
    if buf:
        tables.append("\n".join(buf))
    return tables


# ---------------------------------------------------------------------------
# DOCX — python-docx
# ---------------------------------------------------------------------------

def _parse_docx(path: Path) -> dict[str, Any]:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph
    from docx.document import Document as _Doc

    doc = Document(str(path))
    lines: list[str] = []
    tables_text: list[str] = []
    title: str | None = None

    def iter_block_items(parent) -> Iterator[dict[str, Any]]:
        if isinstance(parent, _Doc):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            parent_elm = parent._element
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield {"type": "p", "obj": Paragraph(child, parent)}
            elif isinstance(child, CT_Tbl):
                yield {"type": "table", "obj": Table(child, parent)}

    for block in iter_block_items(doc):
        if block["type"] == "p":
            para = block["obj"]
            text = para.text.strip()
            if not text:
                lines.append("")
                continue
            style = (para.style.name or "").lower() if para.style else ""
            if title is None and (style.startswith("heading 1") or style == "title"):
                title = text
            if style.startswith("heading 1"):
                lines.append(f"# {text}")
            elif style.startswith("heading 2"):
                lines.append(f"## {text}")
            elif style.startswith("heading 3"):
                lines.append(f"### {text}")
            elif style.startswith("heading 4"):
                lines.append(f"#### {text}")
            elif "list" in style and "number" in style:
                lines.append(f"1. {text}")
            elif "list" in style:
                lines.append(f"- {text}")
            else:
                rich = _para_to_markdown(para)
                lines.append(rich if rich else text)
        elif block["type"] == "table":
            t = block["obj"]
            tbl_md = _table_to_markdown(t)
            lines.append(tbl_md)
            tables_text.append(tbl_md)

    body = "\n".join(lines)
    if title is None:
        title = _first_nonempty_line(body)

    return {
        "title": title,
        "body": body,
        "tables": tables_text,
        "meta": {
            "path": str(path),
            "size_kb": round(path.stat().st_size / 1024, 1),
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
        },
        "source_format": "docx",
    }


def _para_to_markdown(para) -> str:
    parts: list[str] = []
    for run in para.runs:
        t = run.text or ""
        if not t:
            continue
        bold = bool(run.bold)
        italic = bool(run.italic)
        if bold and italic:
            t = f"***{t}***"
        elif bold:
            t = f"**{t}**"
        elif italic:
            t = f"*{t}*"
        parts.append(t)
    return "".join(parts).strip()


def _table_to_markdown(table) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells = [c.text.replace("\n", " ").strip() for c in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if rows:
        ncols = rows[0].count("|") - 1
        if ncols > 0:
            sep = "| " + " | ".join(["---"] * ncols) + " |"
            rows.insert(1, sep)
    return "\n".join(rows)


def _first_nonempty_line(text: str) -> str | None:
    for line in text.splitlines():
        s = line.strip()
        if s:
            return s.lstrip("#").strip()
    return None


# ---------------------------------------------------------------------------
# PDF — pypdf
# ---------------------------------------------------------------------------

def _parse_pdf(path: Path) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages_text: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            t = page.extract_text() or ""
        except Exception as e:
            t = f"\n[Page {i+1}: text extraction failed: {e}]\n"
        pages_text.append(t)
    body = "\n\n".join(pages_text)

    title = None
    if reader.metadata and reader.metadata.title:
        title = str(reader.metadata.title).strip()
    if not title:
        title = _first_nonempty_line(body)

    return {
        "title": title,
        "body":  body,
        "tables": [],
        "meta": {
            "path": str(path),
            "size_kb": round(path.stat().st_size / 1024, 1),
            "page_count": len(reader.pages),
        },
        "source_format": "pdf",
    }


# ---------------------------------------------------------------------------
# PPTX — python-pptx
# ---------------------------------------------------------------------------

def _parse_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    pres = Presentation(str(path))
    lines: list[str] = []
    title: str | None = None

    for i, slide in enumerate(pres.slides, start=1):
        slide_title: str | None = None
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if not txt:
                    continue
                if slide_title is None and para.level == 0:
                    slide_title = txt
                if slide_title == txt and title is None:
                    title = txt
                if slide_title == txt:
                    slide_lines.append(f"### {txt}")
                else:
                    slide_lines.append(f"- {txt}")
        if slide_title is None:
            slide_title = f"Slide {i}"
        lines.append(f"## Slide {i}: {slide_title}")
        lines.extend(slide_lines)
        lines.append("")

    body = "\n".join(lines).rstrip()
    if title is None:
        title = path.stem

    return {
        "title": title,
        "body":  body,
        "tables": [],
        "meta": {
            "path": str(path),
            "size_kb": round(path.stat().st_size / 1024, 1),
            "slide_count": len(pres.slides),
        },
        "source_format": "pptx",
    }


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

def parse(path: str | Path) -> dict[str, Any]:
    """Parse any supported file. Raises ValueError for unsupported formats."""
    p = Path(path)
    ext = p.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(
            f"Unsupported format: {ext}. Supported: {sorted(SUPPORTED_EXTS)}"
        )
    if ext in (".md", ".txt"):
        return _parse_markdown_or_text(p)
    if ext == ".docx":
        return _parse_docx(p)
    if ext == ".pdf":
        return _parse_pdf(p)
    if ext == ".pptx":
        return _parse_pptx(p)
    raise ValueError(f"Unhandled: {ext}")
